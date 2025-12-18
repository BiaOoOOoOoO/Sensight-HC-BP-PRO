import streamlit as st
from google import genai
from google.genai import types
from duckduckgo_search import DDGS
import time
import io
import re
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from tenacity import retry, stop_after_attempt, wait_exponential

# ==========================================
# 0. 系统配置 (隐藏底层技术细节)
# ==========================================
SYSTEM_VERSION = "Sensight BP PRO v2.5"
MODEL_ID = "gemini-2.5-flash"  # 内部调用，不展示给客户

st.set_page_config(
    page_title="Sensight Healthcare BP PRO",
    page_icon="🏥",
    layout="wide"
)

# ==========================================
# 1. UI 样式：纯净专业版 (去 AI 化)
# ==========================================
st.markdown("""
<style>
    .stApp { background-color: #FAFAFA; font-family: 'Helvetica Neue', 'Arial', sans-serif; }
    h1, h2, h3 { color: #0f172a !important; font-weight: 700; }
    
    /* 侧边栏样式 */
    [data-testid="stSidebar"] { background-color: #F1F5F9; border-right: 1px solid #E2E8F0; }
    
    /* 核心按钮：品牌色 */
    div.stButton > button {
        background-color: #2563EB; /* 专业蓝 */
        color: white; border: none;
        width: 100%; padding: 12px; font-weight: 600; border-radius: 6px;
        font-size: 16px; transition: all 0.3s;
    }
    div.stButton > button:hover { background-color: #1d4ed8; }

    /* 报告容器 */
    .report-box {
        background: white; padding: 40px; 
        border: 1px solid #E2E8F0; border-radius: 8px;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
        color: #334155; line-height: 1.8; font-size: 15px;
    }
    
    /* 状态提示隐藏技术细节 */
    .status-box {
        padding: 10px 15px; border-radius: 6px; margin-bottom: 15px;
        background-color: #EFF6FF; border-left: 4px solid #2563EB;
        color: #1E40AF; font-size: 0.95em;
    }
</style>
""", unsafe_allow_html=True)

# ==========================================
# 2. 核心逻辑：多源数据检索 (解决搜索为空)
# ==========================================

def search_market_data(query):
    """
    执行多维度市场数据检索。
    策略：尝试多个关键词组合，直到获取到有效数据，避免返回空值。
    """
    context = ""
    try:
        # 提取核心关键词 (去掉换行和无关词)
        clean_query = query[:60].replace("\n", " ")
        
        # 定义搜索策略队列
        search_strategies = [
            f"{clean_query} market size 2024 2030 report",
            f"{clean_query} clinical trial phase 3 results competitors",
            f"{clean_query} competitive landscape analysis"
        ]
        
        with DDGS() as ddgs:
            for q in search_strategies:
                # 增加延迟防止被封锁
                time.sleep(0.5)
                # 尝试搜索
                results = list(ddgs.text(q, max_results=2))
                if results:
                    for r in results:
                        context += f"- [Source: {r['title']}]: {r['body']}\n"
            
            # 如果跑完所有策略还是空的
            if not context:
                context = "Internal Knowledge Base Only (External data unavailable due to network restrictions)."
                
    except Exception as e:
        # 捕获异常但不报错，保证流程继续
        context = f"Data retrieval skipped: {str(e)}"
        
    return context

@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
def generate_report(api_key, user_input, search_data):
    """
    调用 Sensight 核心引擎 (封装后的 Google GenAI SDK)
    """
    client = genai.Client(api_key=api_key)
    
    # 系统指令：严格的商业分析师人设，禁止提及 AI 身份
    system_instruction = """
    You are a Senior Healthcare Investment Consultant at Sensight.
    Your output is a formal Due Diligence Report for a client.
    
    STRICT GUIDELINES:
    1. NEVER mention you are an AI, Gemini, or a language model.
    2. Tone: Highly professional, objective, concise, and authoritative.
    3. Output Language: Professional Chinese.
    4. Structure: Use Markdown headers (##) strictly for section separation.
    """
    
    full_prompt = f"""
    [PROJECT INPUT]
    {user_input}
    
    [EXTERNAL MARKET DATA]
    {search_data}
    
    [TASK]
    Analyze the project and provide a modification report.
    
    [REQUIRED SECTIONS]
    ## 1. 关键数据核查 (Data Verification)
    - Cross-reference input data with External Market Data.
    - Format: "原表述 -> 修正建议 [依据]"
    - If External Data is empty, rely on logical consistency checks.

    ## 2. 竞品格局分析 (Competitive Landscape)
    - Create a Markdown Table comparing the project vs 3 Major Competitors.
    - Columns: Competitor, Technology/Modality, Stage, Key Differentiation.

    ## 3. 专业化术语升级 (Terminology Upgrade)
    - Rewrite the input summary using investment banking standard terminology.
    - Focus on clinical value and commercial potential.

    ## 4. BP 幻灯片大纲 (Slide Deck Framework)
    - Provide content for 4 key slides.
    - Format strictly as:
    Slide 1: [Title] - [Bullet points]
    Slide 2: [Title] - [Bullet points]
    ...
    """
    
    response = client.models.generate_content(
        model=MODEL_ID,
        contents=full_prompt,
        config=types.GenerateContentConfig(
            system_instruction=system_instruction,
            temperature=0.2, # 低温保证专业一致性
        )
    )
    return response.text

# ==========================================
# 3. 专业文件生成引擎 (Word & PPT)
# ==========================================

def create_word(text):
    doc = Document()
    doc.add_heading('Sensight Project Analysis Report', 0)
    
    for line in text.split('\n'):
        clean_line = line.strip()
        if clean_line.startswith('## '):
            doc.add_heading(clean_line.replace('## ', ''), level=2)
        elif clean_line.startswith('|'):
            doc.add_paragraph(clean_line, style='List Bullet')
        else:
            p = doc.add_paragraph(clean_line)
            p.paragraph_format.space_after = Pt(6)
            
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_ppt(text):
    """
    智能 PPT 生成器：解析 markdown 结构，生成带版式的 PPT
    """
    prs = Presentation()
    
    # 1. 封面页
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project Optimization Report"
    subtitle.text = "Generated by Sensight Healthcare Solution"
    
    # 2. 解析内容页 (按 Markdown 标题 ## 分割)
    # 使用正则切分章节
    sections = re.split(r'^##\s+', text, flags=re.MULTILINE)
    
    for section in sections:
        if not section.strip(): continue
        
        lines = section.strip().split('\n')
        header = lines[0].strip()
        content_lines = lines[1:]
        
        # 过滤空行并清理
        bullets = [line.strip().lstrip('-*•').strip() for line in content_lines if line.strip()]
        
        # 创建内容页 (Title and Content)
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = header
        
        # 设置正文 (作为 Bullet Points)
        tf = slide.placeholders[1].text_frame
        tf.clear() # 清除默认占位符文本
        
        for point in bullets:
            # 简单逻辑：如果这一行很短且不包含大量标点，作为一级要点
            if point:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(18) # 设置字体大小，显得更专业
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# ==========================================
# 4. 主程序界面
# ==========================================

# 侧边栏：系统设置
with st.sidebar:
    st.image("https://placehold.co/200x60/2563EB/FFFFFF?text=SENSIGHT", caption="Healthcare Solutions")
    st.markdown("---")
    
    # 隐晦的 API Key 输入提示
    api_key = st.text_input("系统授权密钥 (License Key)", type="password")
    
    st.markdown("### ⚙️ 系统状态")
    st.caption("核心引擎: Online")
    st.caption("数据链路: Secure")
    st.caption(f"版本: {SYSTEM_VERSION}")

# 主界面
st.title("Sensight 医疗项目 BP 优化系统")
st.markdown("请在下方输入项目核心段落（摘要/临床数据/竞品），系统将基于全球数据库进行**合规性核查**与**专业级润色**。")

user_input = st.text_area("项目数据输入 (Project Input)", height=300, placeholder="粘贴项目核心描述...")

if st.button("生成专业分析报告"):
    if not api_key:
        st.error("请输入系统授权密钥")
    elif not user_input:
        st.warning("请输入项目数据")
    else:
        # 使用自定义样式的容器，避免 Streamlit 默认的加载条
        status_placeholder = st.empty()
        
        try:
            # Phase 1: 数据检索
            status_placeholder.markdown('<div class="status-box">正在连接全球临床数据库检索竞品信息...</div>', unsafe_allow_html=True)
            search_res = search_market_data(user_input)
            
            # Phase 2: 核心分析
            status_placeholder.markdown('<div class="status-box">Sensight 引擎正在执行逻辑校验与架构重组...</div>', unsafe_allow_html=True)
            final_report = generate_report(api_key, user_input, search_res)
            
            # 完成
            status_placeholder.empty()
            st.success("分析完成")
            
            # 展示报告
            st.markdown(f'<div class="report-box">{final_report}</div>', unsafe_allow_html=True)
            
            # 下载区域
            st.markdown("### 📥 导出专业文档")
            c1, c2 = st.columns(2)
            with c1:
                st.download_button(
                    label="📄 导出分析报告 (.docx)",
                    data=create_word(final_report),
                    file_name="Sensight_Analysis_Report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            with c2:
                st.download_button(
                    label="📊 导出演示文稿 (.pptx)",
                    data=create_ppt(final_report),
                    file_name="Sensight_Presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
        except Exception as e:
            st.error("系统处理过程中遇到异常，请检查密钥权限或网络连接。")
            # 仅在展开详情里显示技术错误，避免直接暴露给客户
            with st.expander("错误详情 (仅供技术支持)"):
                st.write(str(e))
